from pathlib import Path
import pandas as pd
from pypdf import PdfReader
import json
from PIL import Image

from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup


class Loader:
    @staticmethod
    def load(path: Path):
        ext = path.suffix.lower()

        # -------------------------------
        # 1. PDF
        # -------------------------------
        if ext == ".pdf":
            reader = PdfReader(str(path))
            text = "\n".join([(p.extract_text() or "") for p in reader.pages])
            return {"type": "text", "content": text}

        # -------------------------------
        # 2. CSV / TSV
        # -------------------------------
        if ext in [".csv", ".tsv"]:
            return {"type": "table", "df": pd.read_csv(path)}

        # -------------------------------
        # 3. Excel
        # -------------------------------
        if ext in [".xls", ".xlsx"]:
            return {"type": "table", "df": pd.read_excel(path)}

        # -------------------------------
        # 4. JSON
        # -------------------------------
        if ext == ".json":
            return {"type": "json", "json": json.loads(path.read_text())}

        # -------------------------------
        # 5. Images
        # -------------------------------
        if ext in [".jpg", ".jpeg", ".png", ".bmp", ".gif"]:
            return {"type": "image", "image": Image.open(path)}

        # -------------------------------
        # 6. TXT / Markdown
        # -------------------------------
        if ext in [".txt", ".md"]:
            return {"type": "text", "content": path.read_text(errors="ignore")}

        # -------------------------------
        # 7. DOCX
        # -------------------------------
        if ext == ".docx":
            doc = DocxDocument(path)
            text = "\n".join([p.text for p in doc.paragraphs])
            return {"type": "text", "content": text}

        # -------------------------------
        # 8. PPTX
        # -------------------------------
        if ext == ".pptx":
            prs = Presentation(path)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            return {"type": "text", "content": "\n".join(slides_text)}

        # -------------------------------
        # 9. HTML / HTM
        # -------------------------------
        if ext in [".html", ".htm"]:
            html = path.read_text(errors="ignore")
            soup = BeautifulSoup(html, "html.parser")
            text = soup.get_text(separator=" ")
            return {"type": "text", "content": text}

        # -------------------------------
        # 10. Default text loader
        # -------------------------------
        return {"type": "text", "content": path.read_text(errors="ignore")}
