# Export Engine - Delivery Intelligence
"""
Generates executive-ready exports from chat responses.
Supports PDF, PPTX, and Email HTML formats.

No user rework required - one-click export with company branding.
"""

import io
import base64
from datetime import datetime
from typing import Dict, List, Optional, Any
from dataclasses import dataclass
import json

# PDF generation
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
    from reportlab.lib.units import inch
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("⚠️ ReportLab not installed - PDF export disabled")

# PPTX generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx not installed - PPTX export disabled")


@dataclass
class ExportContent:
    """Content to export"""
    title: str
    answer: str
    charts: List[Dict] = None  # List of chart data/images
    tables: List[Dict] = None  # List of table data
    metadata: Dict = None      # Additional metadata
    
    def __post_init__(self):
        self.charts = self.charts or []
        self.tables = self.tables or []
        self.metadata = self.metadata or {}


@dataclass
class CompanyBranding:
    """Company branding for exports"""
    company_name: str = "Business Report"
    primary_color: str = "#3B82F6"
    secondary_color: str = "#1E40AF"
    logo_base64: Optional[str] = None
    font_family: str = "Helvetica"


class PDFExporter:
    """Generates executive PDF reports"""
    
    def __init__(self, branding: CompanyBranding = None):
        self.branding = branding or CompanyBranding()
    
    def export(self, content: ExportContent) -> bytes:
        """Generate PDF from content"""
        if not PDF_AVAILABLE:
            raise RuntimeError("ReportLab not installed. Run: pip install reportlab")
        
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.5*inch, bottomMargin=0.5*inch)
        
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            textColor=colors.HexColor(self.branding.primary_color),
            spaceAfter=12,
        )
        
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontSize=10,
            spaceAfter=8,
            leading=14,
        )
        
        # Build story
        story = []
        
        # Header with branding
        story.append(Paragraph(self.branding.company_name, title_style))
        story.append(Paragraph(content.title, styles['Heading2']))
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}", styles['Normal']))
        story.append(Spacer(1, 0.25*inch))
        
        # Main content
        # Split answer into paragraphs and format
        for paragraph in content.answer.split('\n\n'):
            if paragraph.strip():
                # Handle markdown bold
                paragraph = paragraph.replace('**', '<b>').replace('**', '</b>')
                story.append(Paragraph(paragraph.strip(), body_style))
        
        story.append(Spacer(1, 0.5*inch))
        
        # Tables
        for table_data in content.tables:
            if table_data.get('rows'):
                headers = table_data.get('headers', [])
                rows = table_data.get('rows', [])
                
                table_rows = [headers] if headers else []
                table_rows.extend(rows)
                
                if table_rows:
                    t = Table(table_rows)
                    t.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.branding.primary_color)),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 10),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ]))
                    story.append(t)
                    story.append(Spacer(1, 0.25*inch))
        
        # Footer
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(
            f"© {datetime.now().year} {self.branding.company_name} | AI Business Analyst",
            styles['Normal']
        ))
        
        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()


class PPTXExporter:
    """Generates executive PowerPoint presentations"""
    
    def __init__(self, branding: CompanyBranding = None):
        self.branding = branding or CompanyBranding()
    
    def export(self, content: ExportContent) -> bytes:
        """Generate PPTX from content"""
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Add title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Add subtitle with company name
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"{self.branding.company_name} | {datetime.now().strftime('%B %d, %Y')}"
        p2.font.size = Pt(24)
        p2.alignment = PP_ALIGN.CENTER
        
        # Content slide
        content_layout = prs.slide_layouts[6]  # Blank
        slide2 = prs.slides.add_slide(content_layout)
        
        # Add content
        txBox3 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        
        # Parse content (simplified)
        paragraphs = content.answer.split('\n\n')[:5]  # Limit for slide
        for i, para in enumerate(paragraphs):
            if i == 0:
                p = tf3.paragraphs[0]
            else:
                p = tf3.add_paragraph()
            
            # Clean markdown
            para = para.replace('**', '').replace('*', '')
            p.text = para[:300]  # Limit length
            p.font.size = Pt(18 if i == 0 else 14)
            p.font.bold = (i == 0)
        
        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()


class EmailExporter:
    """Generates executive email HTML"""
    
    def __init__(self, branding: CompanyBranding = None):
        self.branding = branding or CompanyBranding()
    
    def export(self, content: ExportContent) -> str:
        """Generate email HTML from content"""
        
        # Convert markdown-ish to HTML
        body_html = self._markdown_to_html(content.answer)
        
        html = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Helvetica, Arial, sans-serif;
            line-height: 1.6;
            color: #333;
            max-width: 600px;
            margin: 0 auto;
            padding: 20px;
        }}
        .header {{
            background: linear-gradient(135deg, {self.branding.primary_color}, {self.branding.secondary_color});
            color: white;
            padding: 30px 20px;
            border-radius: 8px 8px 0 0;
            text-align: center;
        }}
        .header h1 {{
            margin: 0;
            font-size: 24px;
        }}
        .header p {{
            margin: 10px 0 0 0;
            opacity: 0.9;
            font-size: 14px;
        }}
        .content {{
            background: #ffffff;
            border: 1px solid #e0e0e0;
            border-top: none;
            padding: 30px 20px;
            border-radius: 0 0 8px 8px;
        }}
        .content h2 {{
            color: {self.branding.primary_color};
            font-size: 18px;
            margin-top: 0;
        }}
        .content p {{
            margin: 15px 0;
        }}
        .key-insight {{
            background: #f8f9fa;
            border-left: 4px solid {self.branding.primary_color};
            padding: 15px 20px;
            margin: 20px 0;
        }}
        .footer {{
            text-align: center;
            margin-top: 30px;
            padding: 20px;
            color: #666;
            font-size: 12px;
        }}
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 20px 0;
        }}
        th, td {{
            border: 1px solid #ddd;
            padding: 10px;
            text-align: left;
        }}
        th {{
            background: {self.branding.primary_color};
            color: white;
        }}
        tr:nth-child(even) {{
            background: #f9f9f9;
        }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{self.branding.company_name}</h1>
        <p>Business Intelligence Report | {datetime.now().strftime('%B %d, %Y')}</p>
    </div>
    <div class="content">
        <h2>{content.title}</h2>
        {body_html}
    </div>
    <div class="footer">
        Generated by AI Business Analyst<br>
        © {datetime.now().year} {self.branding.company_name}
    </div>
</body>
</html>
"""
        return html
    
    def _markdown_to_html(self, text: str) -> str:
        """Convert markdown-like text to HTML"""
        import re
        
        # Bold
        text = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', text)
        text = re.sub(r'__(.+?)__', r'<strong>\1</strong>', text)
        
        # Italic
        text = re.sub(r'\*(.+?)\*', r'<em>\1</em>', text)
        
        # Convert paragraphs
        paragraphs = text.split('\n\n')
        html_parts = []
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            # Check if it's a key insight (starts with bold)
            if para.startswith('<strong>'):
                html_parts.append(f'<div class="key-insight">{para}</div>')
            else:
                html_parts.append(f'<p>{para}</p>')
        
        return '\n'.join(html_parts)


class ExportEngine:
    """
    Unified export engine for all formats.
    """
    
    def __init__(self, branding: CompanyBranding = None):
        self.branding = branding or CompanyBranding()
        self.pdf = PDFExporter(self.branding)
        self.pptx = PPTXExporter(self.branding)
        self.email = EmailExporter(self.branding)
    
    def export_pdf(self, content: ExportContent) -> bytes:
        """Export to PDF"""
        return self.pdf.export(content)
    
    def export_pptx(self, content: ExportContent) -> bytes:
        """Export to PowerPoint"""
        return self.pptx.export(content)
    
    def export_email(self, content: ExportContent) -> str:
        """Export to HTML email"""
        return self.email.export(content)
    
    def get_available_formats(self) -> List[str]:
        """Get list of available export formats"""
        formats = ['email']  # Email always available (pure HTML)
        if PDF_AVAILABLE:
            formats.append('pdf')
        if PPTX_AVAILABLE:
            formats.append('pptx')
        return formats


def get_export_engine(workspace_id: str = None) -> ExportEngine:
    """
    Get export engine with company branding if available.
    """
    branding = CompanyBranding()
    
    if workspace_id:
        try:
            from core.company_profile import get_company_profile, get_company_branding
            
            profile = get_company_profile(workspace_id)
            if profile:
                company_branding = get_company_branding(workspace_id)
                branding = CompanyBranding(
                    company_name=profile.company_name,
                    primary_color=company_branding.primary_color,
                    secondary_color=company_branding.secondary_color,
                    logo_base64=company_branding.logo_base64,
                    font_family=company_branding.font_family,
                )
        except Exception as e:
            print(f"⚠️ Could not load company branding: {e}")
    
    return ExportEngine(branding)
